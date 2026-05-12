from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .models import Box, ImageElement, PresentationSpec, SlideSpec, TextElement


class RenderingError(Exception):
    pass


@dataclass
class TextRenderer:
    def render(self, slide, element: TextElement) -> None:
        b = element.box
        tx = slide.shapes.add_textbox(Emu(b.x), Emu(b.y), Emu(b.width), Emu(b.height))
        tf = tx.text_frame
        tf.word_wrap = element.style.wrap
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }[element.style.align]
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[element.style.valign]

        run = p.add_run()
        run.text = element.text
        font = run.font
        font.name = element.style.font_name
        font.bold = element.style.bold
        font.italic = element.style.italic
        font.color.rgb = RGBColor.from_string(element.style.color_hex)

        size = element.style.font_size_pt
        if element.style.auto_fit_font:
            size = self._fit_font_size(element)
        font.size = Pt(size)

    def _fit_font_size(self, element: TextElement) -> float:
        area = max(element.box.width * element.box.height, 1)
        char_count = max(len(element.text), 1)
        approx = (area / char_count) ** 0.5 / 1000.0 * 7.2
        return max(element.style.min_font_size_pt, min(element.style.max_font_size_pt, approx))


@dataclass
class ImageRenderer:
    def render(self, slide, element: ImageElement) -> None:
        image_path = Path(element.path)
        if not image_path.exists():
            raise RenderingError(f"Image not found: {image_path}")

        with Image.open(image_path) as img:
            img_w, img_h = img.size

        b = element.box
        if element.crop_mode == "fit":
            slide.shapes.add_picture(str(image_path), Emu(b.x), Emu(b.y), Emu(b.width), Emu(b.height))
            return

        box_ratio = b.width / b.height
        image_ratio = img_w / img_h

        if element.crop_mode == "contain":
            if image_ratio > box_ratio:
                width = b.width
                height = int(width / image_ratio)
                x, y = b.x, b.y + (b.height - height) // 2
            else:
                height = b.height
                width = int(height * image_ratio)
                x, y = b.x + (b.width - width) // 2, b.y
            slide.shapes.add_picture(str(image_path), Emu(x), Emu(y), Emu(width), Emu(height))
            return

        if element.crop_mode == "cover":
            pic = slide.shapes.add_picture(str(image_path), Emu(b.x), Emu(b.y), Emu(b.width), Emu(b.height))
            if image_ratio > box_ratio:
                visible_ratio = box_ratio / image_ratio
                crop_each = (1.0 - visible_ratio) / 2.0
                pic.crop_left = crop_each
                pic.crop_right = crop_each
            else:
                visible_ratio = image_ratio / box_ratio
                crop_each = (1.0 - visible_ratio) / 2.0
                pic.crop_top = crop_each
                pic.crop_bottom = crop_each
            return

        raise RenderingError(f"Unsupported crop mode: {element.crop_mode}")


@dataclass
class SlideRenderer:
    text_renderer: TextRenderer
    image_renderer: ImageRenderer

    def render(self, prs: Presentation, slide_spec: SlideSpec) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for element in slide_spec.elements:
            if isinstance(element, TextElement):
                self.text_renderer.render(slide, element)
            elif isinstance(element, ImageElement):
                self.image_renderer.render(slide, element)
            else:
                raise RenderingError(f"Unsupported element: {type(element)!r}")


@dataclass
class Renderer:
    slide_renderer: SlideRenderer

    def render(self, spec: PresentationSpec, output_path: str) -> None:
        prs = Presentation()
        prs.slide_width = Emu(spec.slide_size.width)
        prs.slide_height = Emu(spec.slide_size.height)

        for slide_spec in spec.slides:
            self.slide_renderer.render(prs, slide_spec)

        prs.save(output_path)
